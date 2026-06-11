# app/services/parser.py - 文件解析器服务
# 功能：解析多种格式文件（PDF, DOCX, PPTX, XLS 等）并提取文本内容

import os
import re
import time
import asyncio
import pymupdf4llm
import docx
from pptx import Presentation
import pandas as pd
import subprocess
from typing import Union, List, Optional
from app.core.exceptions import FileProcessingException
from app.core.logging import logger


def clean_surrogate_characters(text: str) -> str:
    """
    清理文本中的 Unicode 代理字符（surrogate characters）。
    
    代理字符是 Unicode 中 U+D800 到 U+DFFF 范围内的字符，它们用于 UTF-16 编码中的
    辅助平面字符表示，但不能单独存在于有效的 UTF-8/UTF-32 文本中。
    
    PDF 文本提取时可能会产生这些无效字符，导致保存到 MySQL 时出现编码错误。
    
    Args:
        text: 输入文本
        
    Returns:
        清理后的文本
    """
    if not text:
        return text
    
    # 移除所有代理字符（U+D800 到 U+DFFF）
    # 这些字符在 Python 字符串中可能以 '\ud800' 到 '\udfff' 形式出现
    cleaned_text = re.sub(r'[\ud800-\udfff]', '', text)
    
    # 记录是否有字符被清理
    if len(cleaned_text) < len(text):
        removed_count = len(text) - len(cleaned_text)
        logger.warning(f"[clean_surrogate_characters] 已移除 {removed_count} 个 Unicode 代理字符")
    
    return cleaned_text


def _parse_excel(fname: str) -> str:
    """
    高质量 Excel 文件解析，支持 .xlsx 和 .xls 格式。
    
    特性：
    - 使用 openpyxl 直接解析 .xlsx，保留合并单元格、表头等结构信息
    - 使用 xlrd 解析 .xls（旧格式）
    - 输出 Markdown 表格格式，对 LLM 友好
    - 自动处理合并单元格（填充展开值）
    - 智能检测表头行
    - 多 Sheet 分别输出
    - Fallback 到 pandas 解析
    """
    fname_lower = fname.lower()
    excel_parts = []

    # ========== .xlsx 格式：使用 openpyxl 直接解析 ==========
    if fname_lower.endswith('.xlsx'):
        try:
            from openpyxl import load_workbook
            # 注意：read_only=True 模式不支持 merged_cells，因此使用普通模式以保留合并单元格信息
            wb = load_workbook(fname, data_only=True)
            
            for sheet_idx, sheet_name in enumerate(wb.sheetnames):
                ws = wb[sheet_name]
                rows_data = []
                
                # 收集合并单元格信息
                merged_map = {}
                for merged_range in ws.merged_cells.ranges:
                    min_row, min_col, max_row, max_col = (
                        merged_range.min_row, merged_range.min_col,
                        merged_range.max_row, merged_range.max_col
                    )
                    # 获取合并区域左上角的值
                    cell_value = ws.cell(row=min_row, column=min_col).value
                    # 将合并区域的所有单元格映射到左上角的值
                    for r in range(min_row, max_row + 1):
                        for c in range(min_col, max_col + 1):
                            if r != min_row or c != min_col:
                                merged_map[(r, c)] = cell_value
                
                # 逐行读取数据
                for row in ws.iter_rows():
                    row_values = []
                    for cell in row:
                        # 优先使用合并单元格映射的值
                        if (cell.row, cell.column) in merged_map:
                            val = merged_map[(cell.row, cell.column)]
                        else:
                            val = cell.value
                        
                        if val is None:
                            row_values.append("")
                        elif isinstance(val, (int, float)):
                            # 保留数值精度，避免科学计数法
                            if isinstance(val, float) and val == int(val):
                                row_values.append(str(int(val)))
                            else:
                                row_values.append(str(val))
                        else:
                            row_values.append(str(val).strip())
                    
                    # 跳过完全空行
                    if any(v for v in row_values):
                        rows_data.append(row_values)
                
                wb.close()
                
                if not rows_data:
                    continue
                
                # 智能检测表头行：第一行如果包含非数字、非空内容，视为表头
                header_row = rows_data[0] if rows_data else []
                data_rows = rows_data[1:] if len(rows_data) > 1 else []
                
                # 过滤掉全空的数据行
                data_rows = [r for r in data_rows if any(v for v in r)]
                
                if not any(v for v in header_row):
                    # 表头全空，跳过此 Sheet
                    continue
                
                # 构建 Markdown 表格
                col_count = len(header_row)
                md_lines = []
                md_lines.append("| " + " | ".join(header_row) + " |")
                md_lines.append("| " + " | ".join(["---"] * col_count) + " |")
                
                for row in data_rows:
                    # 补齐列数
                    while len(row) < col_count:
                        row.append("")
                    md_lines.append("| " + " | ".join(row[:col_count]) + " |")
                
                sheet_content = f"\n### Sheet: {sheet_name}\n\n" + "\n".join(md_lines) + "\n"
                excel_parts.append(sheet_content)
                logger.info(f"[Excel解析] 已处理 Sheet {sheet_idx+1}/{len(wb.sheetnames)}: {sheet_name}, "
                           f"表头列数={col_count}, 数据行数={len(data_rows)}")
            
            if excel_parts:
                return "\n".join(excel_parts)
            
        except Exception as e:
            logger.warning(f"[Excel解析] openpyxl 解析 .xlsx 失败: {e}，尝试 pandas fallback...")
    
    # ========== .xls 格式：使用 xlrd 解析 ==========
    elif fname_lower.endswith('.xls'):
        try:
            import xlrd
            wb = xlrd.open_workbook(fname)
            
            for sheet_idx in range(wb.nsheets):
                sheet = wb.sheet_by_index(sheet_idx)
                sheet_name = sheet.name
                
                if sheet.nrows == 0 or sheet.ncols == 0:
                    continue
                
                rows_data = []
                for row_idx in range(sheet.nrows):
                    row_values = []
                    for col_idx in range(sheet.ncols):
                        cell = sheet.cell(row_idx, col_idx)
                        if cell.ctype == xlrd.XL_CELL_TEXT:
                            row_values.append(cell.value.strip())
                        elif cell.ctype == xlrd.XL_CELL_NUMBER:
                            val = cell.value
                            if val == int(val):
                                row_values.append(str(int(val)))
                            else:
                                row_values.append(str(val))
                        elif cell.ctype == xlrd.XL_CELL_DATE:
                            # 尝试格式化日期
                            try:
                                date_val = xlrd.xldate_as_tuple(cell.value, wb.datemode)
                                row_values.append(f"{date_val[0]}-{date_val[1]:02d}-{date_val[2]:02d}")
                            except Exception:
                                row_values.append(str(cell.value))
                        elif cell.ctype == xlrd.XL_CELL_BOOLEAN:
                            row_values.append("是" if cell.value else "否")
                        elif cell.ctype == xlrd.XL_CELL_EMPTY:
                            row_values.append("")
                        else:
                            row_values.append(str(cell.value).strip())
                    
                    if any(v for v in row_values):
                        rows_data.append(row_values)
                
                if not rows_data:
                    continue
                
                header_row = rows_data[0]
                data_rows = rows_data[1:]
                data_rows = [r for r in data_rows if any(v for v in r)]
                
                if not any(v for v in header_row):
                    continue
                
                col_count = len(header_row)
                md_lines = []
                md_lines.append("| " + " | ".join(header_row) + " |")
                md_lines.append("| " + " | ".join(["---"] * col_count) + " |")
                
                for row in data_rows:
                    while len(row) < col_count:
                        row.append("")
                    md_lines.append("| " + " | ".join(row[:col_count]) + " |")
                
                sheet_content = f"\n### Sheet: {sheet_name}\n\n" + "\n".join(md_lines) + "\n"
                excel_parts.append(sheet_content)
                logger.info(f"[Excel解析] 已处理 Sheet {sheet_idx+1}/{wb.nsheets}: {sheet_name}, "
                           f"表头列数={col_count}, 数据行数={len(data_rows)}")
            
            if excel_parts:
                return "\n".join(excel_parts)
            
        except ImportError:
            logger.warning("[Excel解析] xlrd 未安装，尝试 pandas fallback...")
        except Exception as e:
            logger.warning(f"[Excel解析] xlrd 解析 .xls 失败: {e}，尝试 pandas fallback...")
    
    # ========== Fallback: 使用 pandas 解析 ==========
    try:
        engine = 'openpyxl' if fname_lower.endswith('.xlsx') else None
        dfs = pd.read_excel(fname, sheet_name=None, engine=engine)
        
        for sheet_idx, (sheet_name, df) in enumerate(dfs.items()):
            if df.empty:
                continue
            
            # 填充空值，转为字符串
            df = df.fillna("").astype(str)
            # 清理 pandas 产生的浮点数显示
            for col in df.columns:
                df[col] = df[col].apply(lambda x: str(int(float(x))) if x.endswith('.0') and x.replace('.', '', 1).replace('-', '', 1).isdigit() else x)
            
            # 构建 Markdown 表格
            header_row = [str(col) for col in df.columns]
            col_count = len(header_row)
            md_lines = []
            md_lines.append("| " + " | ".join(header_row) + " |")
            md_lines.append("| " + " | ".join(["---"] * col_count) + " |")
            
            for _, row in df.iterrows():
                row_values = [str(v).strip() for v in row.values]
                md_lines.append("| " + " | ".join(row_values) + " |")
            
            sheet_content = f"\n### Sheet: {sheet_name}\n\n" + "\n".join(md_lines) + "\n"
            excel_parts.append(sheet_content)
            logger.info(f"[Excel解析] pandas fallback 已处理 Sheet {sheet_idx+1}/{len(dfs)}: {sheet_name}")
        
        return "\n".join(excel_parts)
    
    except Exception as e:
        logger.error(f"[Excel解析] pandas fallback 也失败: {e}")
        return ""


def process_files(file_list: Union[List, str], vl_enabled: bool = False) -> str:
    """
    解析文件列表，提取文本内容。
    支持 PDF, DOCX, DOC, PPTX, XLSX, XLS, TXT, MD, CSV。
    
    Args:
        file_list: 文件路径列表
        vl_enabled: 是否启用 VL 视觉模型解析（针对 PDF/DOCX 中的图片内容）
    """
    if not file_list:
        logger.info("[文件解析] 文件列表为空，跳过解析")
        return ""
    
    if vl_enabled:
        return _process_files_with_vl(file_list)
    
    text_accumulated = ""
    # 兼容 Gradio 的文件列表和普通路径列表
    files = file_list if isinstance(file_list, list) else [file_list]
    
    logger.info(f"[文件解析] 开始解析 {len(files)} 个文件")
    
    for idx, f in enumerate(files):
        parse_start_time = time.time()
        try:
            # 获取文件名和路径
            fname = f.name if hasattr(f, 'name') else str(f)
            base_name = os.path.basename(fname)
            fname_lower = fname.lower()
            file_content = ""
            
            logger.info(f"[文件解析] [{idx+1}/{len(files)}] 开始处理文件：{base_name}")
            
            # 使用 'file' 命令检测实际的文件类型 (MIME type)
            mime_start_time = time.time()
            try:
                mime_result = subprocess.run(['file', '--mime-type', '-b', fname], capture_output=True, text=True, check=True)
                mime_type = mime_result.stdout.strip()
                logger.info(f"[文件解析] [{idx+1}/{len(files)}] 文件类型检测结果：MIME={mime_type}, 耗时={time.time()-mime_start_time:.2f}s")
            except Exception as e:
                logger.warning(f"[文件解析] [{idx+1}/{len(files)}] 文件类型检测失败：{e}")
                mime_type = ""

            # 优先根据 MIME 类型处理，如果检测失败则根据后缀处理
            if mime_type == "application/pdf" or fname_lower.endswith(".pdf"):
                logger.info(f"[文件解析] [{idx+1}/{len(files)}] 识别为 PDF 文件，使用 pymupdf4llm 提取文本...")
                pdf_start_time = time.time()
                
                # 使用 pymupdf4llm 提取 PDF 内容为 Markdown 格式
                # 这种格式更适合 LLM 处理，能保留表格、标题等结构信息
                try:
                    md_text = pymupdf4llm.to_markdown(fname)
                    file_content = md_text
                    logger.info(f"[文件解析] [{idx+1}/{len(files)}] PDF 文本提取完成（pymupdf4llm），总耗时={time.time()-pdf_start_time:.2f}s, 内容长度={len(file_content)} 字符")
                except Exception as pdf_error:
                    logger.warning(f"[文件解析] [{idx+1}/{len(files)}] pymupdf4llm 提取失败：{pdf_error}，尝试使用基础提取...")
                    # Fallback: 使用 pymupdf 基础提取
                    try:
                        import pymupdf
                        doc = pymupdf.open(fname)
                        page_texts = []
                        for page_idx, page in enumerate(doc):
                            page_text = page.get_text() or ""
                            page_texts.append(page_text)
                        file_content = "\n".join(page_texts)
                        doc.close()
                        logger.info(f"[文件解析] [{idx+1}/{len(files)}] PDF 基础提取完成，总耗时={time.time()-pdf_start_time:.2f}s, 内容长度={len(file_content)} 字符")
                    except Exception as fallback_error:
                        logger.error(f"[文件解析] [{idx+1}/{len(files)}] PDF 提取完全失败：{fallback_error}")
                        file_content = ""
            
            elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or fname_lower.endswith(".docx"):
                logger.info(f"[文件解析] [{idx+1}/{len(files)}] 识别为 DOCX 文件，开始提取文本...")
                docx_start_time = time.time()
                doc = docx.Document(fname)
                content_list = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        content_list.append(para.text)
                for table in doc.tables:
                    for row in table.rows:
                        row_cells = [cell.text.strip() for cell in row.cells]
                        if any(row_cells):
                            content_list.append(" | ".join(row_cells))
                file_content = "\n".join(content_list)
                logger.info(f"[文件解析] [{idx+1}/{len(files)}] DOCX 文本提取完成，耗时={time.time()-docx_start_time:.2f}s, 内容长度={len(file_content)} 字符")

            elif mime_type == "application/msword" or fname_lower.endswith(".doc"):
                # 尝试 antiword
                try:
                    result = subprocess.run(['antiword', fname], capture_output=True, text=True, check=True)
                    file_content = result.stdout
                except:
                    # 如果 antiword 失败，尝试作为 docx 处理
                    try:
                        doc = docx.Document(fname)
                        file_content = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
                    except:
                        # 尝试作为 MHTML 或 HTML 处理 (应对"另存为网页"的 doc)
                        try:
                            with open(fname, "r", encoding="utf-8", errors="ignore") as fo:
                                raw_data = fo.read()
                            
                            if "MIME-Version" in raw_data and "multipart/related" in raw_data:
                                # 解析 MHTML
                                import email
                                from email import policy
                                msg = email.message_from_string(raw_data, policy=policy.default)
                                html_part = ""
                                for part in msg.walk():
                                    if part.get_content_type() == "text/html":
                                        payload = part.get_payload(decode=True)
                                        charset = part.get_content_charset()
                                        # 修复 Word 生成的 MHTML 中常见的 charset=3D"utf-8" 导致 charset 被识别为 "3d" 的问题
                                        if charset and charset.lower().startswith('3d'):
                                            charset = charset[2:].strip('"')
                                        try:
                                            html_part = payload.decode(charset or 'utf-8', errors='ignore')
                                        except:
                                            html_part = payload.decode('utf-8', errors='ignore')
                                        break
                                if html_part:
                                    import re
                                    import html
                                    # 1. 移除 <style> 和 <script> 块
                                    html_part = re.sub(r'<(style|script)[^>]*>.*?</\1>', '', html_part, flags=re.DOTALL | re.IGNORECASE)
                                    # 2. 移除所有 HTML 标签
                                    file_content = re.sub(r'<[^>]+>', '', html_part)
                                    # 3. 还原 HTML 实体 (包括 &#25237; 这种数字编码)
                                    file_content = html.unescape(file_content)
                                    # 4. 还原 Unicode 转义序列 (如 \u3001)
                                    if "\\u" in file_content:
                                        import re
                                        file_content = re.sub(r'\\u([0-9a-fA-F]{4})', lambda m: chr(int(m.group(1), 16)), file_content)
                                    # 5. 清洗空白符
                                    file_content = re.sub(r'\s+', ' ', file_content).strip()
                            elif "<html>" in raw_data.lower():
                                import re
                                import html
                                raw_data = re.sub(r'<(style|script)[^>]*>.*?</\1>', '', raw_data, flags=re.DOTALL | re.IGNORECASE)
                                file_content = re.sub(r'<[^>]+>', '', raw_data)
                                file_content = html.unescape(file_content)
                                if "\\u" in file_content:
                                    file_content = re.sub(r'\\u([0-9a-fA-F]{4})', lambda m: chr(int(m.group(1), 16)), file_content)
                                file_content = re.sub(r'\s+', ' ', file_content).strip()
                            else:
                                # 最后的保底：纯文本
                                file_content = raw_data if len(raw_data) > 10 else ""
                        except:
                            file_content = ""

            elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or fname_lower.endswith(".pptx"):
                logger.info(f"[文件解析] [{idx+1}/{len(files)}] 识别为 PPTX 文件，开始提取文本...")
                pptx_start_time = time.time()
                prs = Presentation(fname)
                ppt_text = []
                for i, slide in enumerate(prs.slides):
                    slide_content = []
                    for shape in slide.shapes:
                        if shape.has_table:
                            table_str = []
                            for row in shape.table.rows:
                                row_cells = [cell.text_frame.text.strip() for cell in row.cells if cell.text_frame]
                                if any(row_cells):
                                    table_str.append(" | ".join(row_cells))
                            if table_str:
                                slide_content.append("\n[表格内容]:\n" + "\n".join(table_str) + "\n")
                        elif shape.has_text_frame:
                            if shape.text_frame.text.strip():
                                slide_content.append(shape.text_frame.text.strip())
                    if slide_content:
                        ppt_text.append(f"[Page {i + 1}]:\n" + "\n".join(slide_content))
                file_content = "\n".join(ppt_text)
                logger.info(f"[文件解析] [{idx+1}/{len(files)}] PPTX 文本提取完成，耗时={time.time()-pptx_start_time:.2f}s, 内容长度={len(file_content)} 字符")

            elif "excel" in mime_type or "spreadsheet" in mime_type or fname_lower.endswith((".xlsx", ".xls")):
                logger.info(f"[文件解析] [{idx+1}/{len(files)}] 识别为 Excel 文件，开始提取文本...")
                excel_start_time = time.time()
                try:
                    file_content = _parse_excel(fname)
                    logger.info(f"[文件解析] [{idx+1}/{len(files)}] Excel 文本提取完成，耗时={time.time()-excel_start_time:.2f}s, 内容长度={len(file_content)} 字符")
                except Exception as e:
                    logger.warning(f"Excel 解析出错 ({base_name}): {str(e)}")
                    file_content = ""
            
            elif "text" in mime_type or fname_lower.endswith((".txt", ".md", ".csv")):
                with open(fname, "r", encoding="utf-8", errors="ignore") as fo:
                    file_content = fo.read()
            
            if file_content.strip():
                # 清理 Unicode 代理字符，避免保存到数据库时出现编码错误
                file_content = clean_surrogate_characters(file_content)
                text_accumulated += f"\n\n<<<<<< 文件开始：{base_name} >>>>>>\n{file_content}\n<<<<<< 文件结束：{base_name} >>>>>>\n"
                logger.info(f"[文件解析] [{idx+1}/{len(files)}] 文件处理成功，内容长度={len(file_content)} 字符")
            else:
                logger.warning(f"[解析警告] 无法从 {base_name} 提取有效内容 (MIME: {mime_type})")
        
        except Exception as e:
            logger.error(f"[读取失败] {base_name}: {str(e)}", exc_info=True)
        finally:
            logger.info(f"[文件解析] [{idx+1}/{len(files)}] 文件处理完成，总耗时={time.time()-parse_start_time:.2f}s")
            
    logger.info(f"[文件解析] 所有文件处理完成，累计内容长度={len(text_accumulated)} 字符")
    return text_accumulated


class FileParser:
    """
    文件解析类，负责多种文件格式的加载与文本提取
    """
    def __init__(self, vl_enabled: bool = False):
        self.vl_enabled = vl_enabled

    def parse_file(self, file_path: str) -> str:
        """
        解析单个文件并返回其文本内容
        """
        return process_files(file_path, vl_enabled=self.vl_enabled)

    def parse_files(self, file_paths: Union[List[str], str]) -> str:
        """
        解析多个文件并合并其文本内容
        """
        return process_files(file_paths, vl_enabled=self.vl_enabled)

    async def async_parse_file(self, file_path: str) -> str:
        """
        异步解析单个文件（VL 模式下使用 VL 解析）
        """
        if self.vl_enabled:
            from app.services.vl_parser import parse_file_with_vl
            loop = asyncio.get_event_loop()
            return await loop.run_in_executor(None, parse_file_with_vl, file_path)
        return process_files(file_path, vl_enabled=False)

    async def async_parse_files(self, file_paths: Union[List[str], str]) -> str:
        """
        异步解析多个文件（VL 模式下使用 VL 解析）
        """
        if self.vl_enabled:
            from app.services.vl_parser import parse_file_with_vl
            files = file_paths if isinstance(file_paths, list) else [file_paths]
            loop = asyncio.get_event_loop()
            results = await asyncio.gather(*[
                loop.run_in_executor(None, parse_file_with_vl, f) for f in files
            ])
            parts = []
            for result, f in zip(results, files):
                if result.strip():
                    base_name = os.path.basename(f)
                    cleaned = clean_surrogate_characters(result)
                    parts.append(f"\n\n<<<<<< 文件开始：{base_name} >>>>>>\n{cleaned}\n<<<<<< 文件结束：{base_name} >>>>>>\n")
            return "".join(parts)
        return process_files(file_paths, vl_enabled=False)


def _process_files_with_vl(file_list: Union[List, str]) -> str:
    """
    使用 VL 视觉模型解析文件列表。
    对于 PDF/DOCX 文件，使用视觉模型识别页面内容（包括图片）；
    对于其他文件类型，回退到传统文本解析。
    """
    from app.services.vl_parser import parse_file_with_vl

    files = file_list if isinstance(file_list, list) else [file_list]
    logger.info(f"[VL解析] 开始处理 {len(files)} 个文件")

    vl_tasks = []
    non_vl_files = []
    file_order = []

    for idx, f in enumerate(files):
        fname = f.name if hasattr(f, 'name') else str(f)
        fname_lower = fname.lower()
        try:
            mime_result = subprocess.run(['file', '--mime-type', '-b', fname], capture_output=True, text=True, check=True)
            mime_type = mime_result.stdout.strip()
        except Exception:
            mime_type = ""

        is_visual = (
            mime_type == "application/pdf" or fname_lower.endswith(".pdf")
            or mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            or fname_lower.endswith(".docx")
        )

        if is_visual:
            vl_tasks.append((idx, fname))
            file_order.append(("vl", idx))
        else:
            non_vl_files.append((idx, fname, f))
            file_order.append(("text", idx))

    vl_results = {}

    if vl_tasks:
        for idx, fname in vl_tasks:
            try:
                result = parse_file_with_vl(fname)
                vl_results[idx] = result
            except Exception as e:
                logger.error(f"[VL解析] 文件 {os.path.basename(fname)} 失败: {e}")
                vl_results[idx] = ""

    text_results = {}
    if non_vl_files:
        for idx, fname, f in non_vl_files:
            text_results[idx] = ""

        non_vl_paths = [f for _, _, f in non_vl_files]
        text_content = process_files(non_vl_paths, vl_enabled=False)

        if text_content.strip():
            file_blocks = text_content.split("<<<<<< 文件开始：")
            for block in file_blocks:
                if not block.strip():
                    continue
                for idx, fname, _ in non_vl_files:
                    base_name = os.path.basename(fname)
                    if block.startswith(base_name):
                        end_marker = f"<<<<<< 文件结束：{base_name} >>>>>>"
                        content = block.replace(f"{base_name} >>>>>>\n", "", 1).replace(end_marker, "").strip()
                        text_results[idx] = content
                        break

    ordered_parts = []
    for file_type, idx in file_order:
        if file_type == "vl":
            content = vl_results.get(idx, "")
        else:
            content = text_results.get(idx, "")

        if content.strip():
            fname = files[idx].name if hasattr(files[idx], 'name') else str(files[idx])
            base_name = os.path.basename(fname)
            content = clean_surrogate_characters(content)
            ordered_parts.append(f"\n\n<<<<<< 文件开始：{base_name} >>>>>>\n{content}\n<<<<<< 文件结束：{base_name} >>>>>>\n")

    result = "".join(ordered_parts)
    logger.info(f"[VL解析] 所有文件处理完成，累计内容长度={len(result)} 字符")
    return result